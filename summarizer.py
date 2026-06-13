import os
import json
import time
import threading
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from datetime import datetime
from openai import OpenAI

# ========================= CONFIG =========================
SOURCE_ROOT      = r"\\bigbertha\Data\Companies_Clients"
RAG_ROOT         = r"D:\RAG"
TRACKER_FILE     = r"D:\RAG\_summary_tracker.json"
BATCH_LIMIT      = 0      # 0 = no limit
MAX_TEXT_CHARS   = 12000  # chars sent to LLM per document
SCAN_MAX_FOLDERS = 0      # 0 = all folders
NUM_WORKERS      = 5      # parallel summarizer threads

LLM_MODEL        = "qwen2.5:1.5b"
LLM_RETRIES      = 5
LLM_RETRY_WAIT   = 15    # seconds between retries

SUPPORTED_EXT    = {'.pdf', '.txt', '.md', '.docx', '.doc', '.pptx', '.ppt',
                    '.xlsx', '.xls', '.csv', '.html', '.htm'}
# =========================================================

llm = OpenAI(api_key="ollama", base_url="http://localhost:11434/v1", timeout=120.0)

_print_lock   = threading.Lock()
_tracker_lock = threading.Lock()


def ts():
    return datetime.now().strftime("%H:%M:%S")


def log(tag: str, event: str, detail: str):
    with _print_lock:
        print(f"[{ts()}] {tag:>7}  {event:<6}  {detail}")


def load_tracker() -> dict:
    try:
        with open(TRACKER_FILE, "r", encoding="utf-8") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {}


def save_tracker(tracker: dict):
    Path(TRACKER_FILE).parent.mkdir(parents=True, exist_ok=True)
    with open(TRACKER_FILE, "w", encoding="utf-8") as f:
        json.dump(tracker, f, indent=2)


def extract_text(file_path: Path) -> tuple[str, int, str]:
    """Return (text, pages_sampled, description). Truncates to MAX_TEXT_CHARS."""
    ext = file_path.suffix.lower()
    sampled = 0

    try:
        if ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            total = len(reader.pages)
            text = ""
            for i, page in enumerate(reader.pages):
                text += (page.extract_text() or "") + "\n"
                sampled = i + 1
                if len(text) >= MAX_TEXT_CHARS:
                    break
            return text[:MAX_TEXT_CHARS], sampled, f"{sampled}/{total} pages"

        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(str(file_path))
            text = "\n".join(p.text for p in doc.paragraphs)
            sampled = len(doc.paragraphs)
            return text[:MAX_TEXT_CHARS], sampled, f"{sampled} paragraphs"

        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(str(file_path))
            total_slides = len(prs.slides)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                sampled += 1
                if len(text) >= MAX_TEXT_CHARS:
                    break
            return text[:MAX_TEXT_CHARS], sampled, f"{sampled}/{total_slides} slides"

        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                text += f"[Sheet: {sheet.title}]\n"
                for row in sheet.iter_rows(values_only=True):
                    row_str = "  |  ".join(str(c) for c in row if c is not None)
                    if row_str.strip():
                        text += row_str + "\n"
                    sampled += 1
                    if len(text) >= MAX_TEXT_CHARS:
                        break
                if len(text) >= MAX_TEXT_CHARS:
                    break
            return text[:MAX_TEXT_CHARS], sampled, f"{sampled} rows"

        elif ext == '.csv':
            text = file_path.read_text(errors='ignore')
            lines = text.splitlines()
            sampled = len(lines)
            return text[:MAX_TEXT_CHARS], sampled, f"{sampled} rows"

        else:  # .txt .md .html .htm
            text = file_path.read_text(errors='ignore')
            sampled = text.count('\n')
            return text[:MAX_TEXT_CHARS], sampled, f"{sampled} lines"

    except Exception as e:
        return f"[Text extraction failed: {e}]", 0, "failed"


def summarize(file_name: str, text: str) -> dict:
    prompt = f"""Analyze the following document thoroughly and return a JSON object with these fields:
- "summary": A detailed summary of 2-3 pages (600-900 words). Cover the main purpose, key findings, arguments, data, conclusions, and any important details. If the document is rich in content expand further so nothing meaningful is lost — the goal is that an LLM reading only this summary can answer detailed questions about the document.
- "key_topics": array of 3-10 key topics or themes
- "document_type": type of document (report, contract, presentation, spreadsheet, email, invoice, etc.)
- "entities": array of key organizations, people, or locations mentioned (max 15)

Respond with ONLY the JSON object — no markdown, no explanation.

Filename: {file_name}
Document text:
{text}"""

    last_err = None
    for attempt in range(1, LLM_RETRIES + 1):
        try:
            response = llm.chat.completions.create(
                model=LLM_MODEL,
                messages=[{"role": "user", "content": prompt}]
            )
            raw = response.choices[0].message.content.strip()
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"):
                    raw = raw[4:]
            try:
                return json.loads(raw)
            except json.JSONDecodeError:
                return {"summary": raw, "key_topics": [], "document_type": "unknown", "entities": []}

        except Exception as e:
            last_err = e
            if attempt < LLM_RETRIES:
                with _print_lock:
                    print(f"[{ts()}] RETRY  attempt {attempt}/{LLM_RETRIES} — {type(e).__name__} — waiting {LLM_RETRY_WAIT}s")
                time.sleep(LLM_RETRY_WAIT)
            else:
                raise RuntimeError(f"LLM failed after {LLM_RETRIES} attempts: {last_err}") from last_err


def rag_path_for(source_path: Path) -> Path:
    rel = source_path.relative_to(SOURCE_ROOT)
    return Path(RAG_ROOT) / rel.parent / (rel.name + ".summary.json")


def process_file(fp: Path, idx: int, total: int) -> tuple[str, str, dict]:
    """Process one file. Returns (file_key, status, tracker_entry)."""
    tag = f"{idx}/{total}"
    size_kb = fp.stat().st_size / 1024
    file_start = time.time()

    log(tag, "START", f"{fp.name}  ({size_kb:,.1f} KB)")

    try:
        # Extract text
        t0 = time.time()
        text, sampled, extract_desc = extract_text(fp)
        extract_secs = time.time() - t0

        if not text.strip() or text.startswith("[Text extraction failed"):
            raise ValueError(f"Extraction failed: {text[:80]}")

        log(tag, "TEXT", f"{len(text):,} chars  {extract_desc}  {extract_secs:.1f}s  — {fp.name}")

        # Summarize
        t0 = time.time()
        analysis = summarize(fp.name, text)
        llm_secs = time.time() - t0

        doc_type  = analysis.get("document_type", "?")
        topics_str = ", ".join(analysis.get("key_topics", [])[:3])
        log(tag, "LLM", f"{llm_secs:.1f}s  type={doc_type}  [{topics_str}]  — {fp.name}")

        # Write summary JSON
        rag_path = rag_path_for(fp)
        summary_obj = {
            "original_path"  : str(fp),
            "rag_path"       : str(rag_path),
            "file_name"      : fp.name,
            "file_extension" : fp.suffix.lower(),
            "file_size_bytes": fp.stat().st_size,
            "file_modified"  : datetime.fromtimestamp(fp.stat().st_mtime).isoformat(),
            "processed_at"   : datetime.now().isoformat(),
            "text_chars_used": len(text),
            "pages_sampled"  : sampled,
            "summary"        : analysis.get("summary", ""),
            "key_topics"     : analysis.get("key_topics", []),
            "document_type"  : doc_type,
            "entities"       : analysis.get("entities", []),
        }
        rag_path.parent.mkdir(parents=True, exist_ok=True)
        with open(rag_path, "w", encoding="utf-8") as f:
            json.dump(summary_obj, f, indent=2, ensure_ascii=False)

        total_secs = time.time() - file_start
        log(tag, "DONE", f"{total_secs:.1f}s total  — {fp.name}")

        return str(fp), "ok", {"status": "ok", "rag_path": str(rag_path), "processed_at": time.time()}

    except Exception as e:
        total_secs = time.time() - file_start
        log(tag, "ERROR", f"{type(e).__name__}: {e}  — {fp.name}")
        return str(fp), "error", {"status": "error", "error": str(e), "processed_at": time.time()}


def main():
    print(f"[{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}] Starting summarizer")
    print(f"  Source  : {SOURCE_ROOT}")
    print(f"  RAG     : {RAG_ROOT}")
    print(f"  Limit   : {BATCH_LIMIT} documents")
    print(f"  Workers : {NUM_WORKERS}\n")

    Path(RAG_ROOT).mkdir(parents=True, exist_ok=True)
    tracker = load_tracker()
    print(f"[{ts()}] Tracker: {len(tracker)} existing entries\n")

    # Determine folders to scan
    source_path = Path(SOURCE_ROOT)
    top_folders = [f for f in source_path.iterdir() if f.is_dir()]
    all_count = len(top_folders)
    if SCAN_MAX_FOLDERS and all_count > SCAN_MAX_FOLDERS:
        top_folders = top_folders[:SCAN_MAX_FOLDERS]
        print(f"[{ts()}] Scanning first {SCAN_MAX_FOLDERS} of {all_count:,} top-level folders")
    else:
        print(f"[{ts()}] Scanning all {all_count:,} top-level folders")
    print()

    # Collect pending — print progress per folder so large shares don't look frozen
    print(f"[{ts()}] Collecting unprocessed files ({len(top_folders):,} folders to scan)...")
    pending = []
    last_print = time.time()
    for fi, folder in enumerate(top_folders, 1):
        folder_count = 0
        for fp in folder.rglob("*"):          # no sorted() — avoids loading whole tree first
            if fp.is_file() and fp.suffix.lower() in SUPPORTED_EXT:
                if str(fp) not in tracker:
                    pending.append(fp)
                    folder_count += 1
        now = time.time()
        if now - last_print >= 5 or fi == len(top_folders):
            print(f"[{ts()}] Scanned {fi:,}/{len(top_folders):,} folders — {len(pending):,} pending so far")
            last_print = now
    if BATCH_LIMIT:
        pending = pending[:BATCH_LIMIT]

    if not pending:
        print("No new documents to summarize.")
        return

    total = len(pending)
    print(f"[{ts()}] Found {total:,} document(s) — launching {NUM_WORKERS} workers\n")
    print("=" * 75)
    print(f"  {'TAG':>7}  {'EVENT':<6}  DETAIL")
    print("=" * 75)

    done = errors = 0
    session_start = time.time()

    with ThreadPoolExecutor(max_workers=NUM_WORKERS) as pool:
        futures = {pool.submit(process_file, fp, i, total): fp
                   for i, fp in enumerate(pending, 1)}

        for future in as_completed(futures):
            file_key, status, entry = future.result()
            with _tracker_lock:
                tracker[file_key] = entry
                save_tracker(tracker)
            if status == "ok":
                done += 1
            else:
                errors += 1

    elapsed = time.time() - session_start
    avg = elapsed / max(done + errors, 1)

    print("\n" + "=" * 75)
    print(f"[{ts()}] Session complete")
    print(f"  Summarized : {done}")
    print(f"  Errors     : {errors}")
    print(f"  Elapsed    : {elapsed:.0f}s  (avg {avg:.0f}s/doc)")
    print(f"  Tracker    : {TRACKER_FILE}")
    print("=" * 75)


if __name__ == "__main__":
    main()
