"""
RAG Pipeline — runs summarizer + ingester in parallel with a live display.
Ctrl+C to stop both.
"""
import os
import json
import time
import threading
from pathlib import Path
from datetime import datetime
from openai import OpenAI

# ========================= CONFIG =========================
SOURCE_ROOT      = r"\\bigbertha\Data\Companies_Clients"
RAG_ROOT         = r"D:\RAG"
SUMMARY_TRACKER  = r"D:\RAG\_summary_tracker.json"
INGEST_TRACKER   = r"D:\RAG\_ingest_tracker.json"

# Summarizer
SUM_MODEL        = "qwen2.5:1.5b"
SUM_MAX_CHARS    = 12000
SUM_WORKERS      = 5
LLM_RETRIES      = 5
LLM_RETRY_WAIT   = 15

# Ingester
ING_COLLECTION   = "rag_summaries"
ING_EMBED_MODEL  = "nomic-embed-text"
ING_CHUNK_SIZE   = 1024
ING_CHUNK_OVERLAP= 64
ING_BATCH_SIZE   = 64
ING_POLL_SECS    = 60    # re-scan for new summaries this often

QDRANT_HOST      = "localhost"
QDRANT_PORT      = 6333

SUPPORTED_EXT    = {'.pdf', '.txt', '.md', '.docx', '.doc', '.pptx', '.ppt',
                    '.xlsx', '.xls', '.csv', '.html', '.htm'}

DISPLAY_REFRESH  = 2
LOG_MAX          = 8
# =========================================================

ollama_llm = OpenAI(api_key="ollama", base_url="http://localhost:11434/v1", timeout=120.0)

# ── Shared state ───────────────────────────────────────────────────────────────
_lock = threading.Lock()

_state = {
    "sum": {
        "status" : "starting",   # starting | scanning | running | done | stopped
        "current": "",
        "step"   : "",
        "done"   : 0,
        "errors" : 0,
        "total"  : 0,            # 0 = not yet known
        "scanned": False,
    },
    "ing": {
        "status"  : "starting",
        "current" : "",
        "step"    : "",
        "done"    : 0,
        "errors"  : 0,
        "vectors" : 0,
        "available": 0,          # summaries available to ingest
    },
    "log": [],   # [(ts, tag, msg), ...]
}

def st(section, **kw):
    with _lock:
        _state[section].update(kw)

def snap(section):
    with _lock:
        return dict(_state[section])

def add_log(tag, msg):
    with _lock:
        _state["log"].append((datetime.now().strftime("%H:%M:%S"), tag, msg))
        if len(_state["log"]) > LOG_MAX:
            _state["log"].pop(0)

def ts():
    return datetime.now().strftime("%H:%M:%S")

# ── Tracker helpers ────────────────────────────────────────────────────────────
def _load_json(path):
    for _ in range(3):
        try:
            with open(path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception:
            time.sleep(0.1)
    return {}

def _save_json(path, data):
    Path(path).parent.mkdir(parents=True, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)

# ══════════════════════════════════════════════════════════════════════════════
# SUMMARIZER
# ══════════════════════════════════════════════════════════════════════════════

def _extract_text(fp: Path) -> tuple[str, int, str]:
    ext = fp.suffix.lower()
    try:
        if ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(str(fp))
            total = len(reader.pages)
            text = ""
            for i, page in enumerate(reader.pages):
                text += (page.extract_text() or "") + "\n"
                if len(text) >= SUM_MAX_CHARS:
                    break
            return text[:SUM_MAX_CHARS], i + 1, f"{i+1}/{total} pages"
        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(str(fp))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text[:SUM_MAX_CHARS], len(doc.paragraphs), f"{len(doc.paragraphs)} paragraphs"
        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(str(fp))
            text = ""
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) >= SUM_MAX_CHARS:
                    break
            return text[:SUM_MAX_CHARS], i + 1, f"{i+1}/{len(prs.slides)} slides"
        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(str(fp), read_only=True, data_only=True)
            text = ""; rows = 0
            for sheet in wb.worksheets:
                text += f"[Sheet: {sheet.title}]\n"
                for row in sheet.iter_rows(values_only=True):
                    r = "  |  ".join(str(c) for c in row if c is not None)
                    if r.strip(): text += r + "\n"
                    rows += 1
                    if len(text) >= SUM_MAX_CHARS: break
                if len(text) >= SUM_MAX_CHARS: break
            return text[:SUM_MAX_CHARS], rows, f"{rows} rows"
        else:
            text = fp.read_text(errors='ignore')
            return text[:SUM_MAX_CHARS], text.count('\n'), f"{text.count(chr(10))} lines"
    except Exception as e:
        return f"[failed: {e}]", 0, "failed"


def _call_llm(file_name: str, text: str) -> dict:
    prompt = f"""Analyze this document and return a JSON object with:
- "summary": detailed 600-900 word summary covering purpose, findings, data, conclusions
- "key_topics": array of 3-10 topics
- "document_type": type of document
- "entities": array of key organizations/people/locations (max 15)

Respond with ONLY the JSON object.

Filename: {file_name}
Text:
{text}"""

    last_err = None
    for attempt in range(1, LLM_RETRIES + 1):
        try:
            r = ollama_llm.chat.completions.create(
                model=SUM_MODEL, messages=[{"role": "user", "content": prompt}])
            raw = r.choices[0].message.content.strip()
            if raw.startswith("```"):
                raw = raw.split("```")[1]
                if raw.startswith("json"): raw = raw[4:]
            try: return json.loads(raw)
            except: return {"summary": raw, "key_topics": [], "document_type": "unknown", "entities": []}
        except Exception as e:
            last_err = e
            if attempt < LLM_RETRIES:
                add_log("SUM", f"LLM retry {attempt}/{LLM_RETRIES} — {type(e).__name__}")
                time.sleep(LLM_RETRY_WAIT)
    raise RuntimeError(f"LLM failed after {LLM_RETRIES} attempts: {last_err}")


def _sum_process_one(fp: Path, tracker: dict) -> str:
    """Process one file for summarization. Returns 'ok' or 'error'."""
    rel = str(fp.relative_to(SOURCE_ROOT))
    name = fp.name

    st("sum", current=name, step="extracting text")
    t0 = time.time()
    text, sampled, desc = _extract_text(fp)

    if not text.strip() or text.startswith("[failed"):
        st("sum", step=f"extraction failed")
        tracker[str(fp)] = {"status": "error", "error": f"extraction failed: {text[:80]}", "processed_at": time.time()}
        add_log("SUM", f"ERR  {name} — extraction failed")
        return "error"

    st("sum", step=f"→ {SUM_MODEL} ({len(text):,} chars)")
    analysis = _call_llm(name, text)
    elapsed = time.time() - t0

    rag_path = Path(RAG_ROOT) / Path(rel).parent / (Path(rel).name + ".summary.json")
    st("sum", step="writing JSON")
    rag_path.parent.mkdir(parents=True, exist_ok=True)
    with open(rag_path, "w", encoding="utf-8") as f:
        json.dump({
            "original_path"  : str(fp),
            "rag_path"       : str(rag_path),
            "file_name"      : name,
            "file_extension" : fp.suffix.lower(),
            "file_size_bytes": fp.stat().st_size,
            "file_modified"  : datetime.fromtimestamp(fp.stat().st_mtime).isoformat(),
            "processed_at"   : datetime.now().isoformat(),
            "text_chars_used": len(text),
            "pages_sampled"  : sampled,
            "summary"        : analysis.get("summary", ""),
            "key_topics"     : analysis.get("key_topics", []),
            "document_type"  : analysis.get("document_type", "unknown"),
            "entities"       : analysis.get("entities", []),
        }, f, indent=2, ensure_ascii=False)

    tracker[str(fp)] = {"status": "ok", "rag_path": str(rag_path), "processed_at": time.time()}
    doc_type = analysis.get("document_type", "?")
    add_log("SUM", f"DONE {name}  {elapsed:.0f}s  [{doc_type}]")
    return "ok"


def summarizer_thread():
    from concurrent.futures import ThreadPoolExecutor, as_completed
    st("sum", status="scanning")
    add_log("SUM", "Scanning source folder...")

    # Scan for pending files
    top_folders = [f for f in Path(SOURCE_ROOT).iterdir() if f.is_dir()]
    st("sum", step=f"scanning {len(top_folders):,} folders...")

    tracker = _load_json(SUMMARY_TRACKER)
    pending = []
    for i, folder in enumerate(top_folders):
        for fp in folder.rglob("*"):
            if fp.is_file() and fp.suffix.lower() in SUPPORTED_EXT:
                if str(fp) not in tracker:
                    pending.append(fp)
        if i % 50 == 0:
            st("sum", step=f"scanning... {i:,}/{len(top_folders):,} folders, {len(pending):,} pending")

    total = len(pending)
    st("sum", status="running", total=total, scanned=True,
       step=f"found {total:,} files")
    add_log("SUM", f"Found {total:,} files to summarize")

    if not pending:
        st("sum", status="done", current="", step="nothing to do")
        return

    tracker_lock = threading.Lock()

    def process(fp):
        result = _sum_process_one(fp, tracker)
        with tracker_lock:
            _save_json(SUMMARY_TRACKER, tracker)
        with _lock:
            if result == "ok":
                _state["sum"]["done"] += 1
            else:
                _state["sum"]["errors"] += 1

    with ThreadPoolExecutor(max_workers=SUM_WORKERS) as pool:
        futures = [pool.submit(process, fp) for fp in pending]
        for _ in as_completed(futures):
            pass

    st("sum", status="done", current="", step=f"complete — {snap('sum')['done']:,} summarized")
    add_log("SUM", "All done")


# ══════════════════════════════════════════════════════════════════════════════
# INGESTER
# ══════════════════════════════════════════════════════════════════════════════

def _build_doc_text(obj: dict) -> str:
    parts = []
    if obj.get("document_type"): parts.append(f"Document type: {obj['document_type']}")
    if obj.get("file_name"):     parts.append(f"File: {obj['file_name']}")
    topics = obj.get("key_topics", [])
    if topics: parts.append(f"Key topics: {', '.join(topics)}")
    entities = obj.get("entities", [])
    if entities: parts.append(f"Entities: {', '.join(entities)}")
    if obj.get("summary"): parts.append(obj["summary"])
    return "\n\n".join(parts)


def ingester_thread():
    from llama_index.core import Settings, Document
    from llama_index.core.ingestion import IngestionPipeline
    from llama_index.core.node_parser import SentenceSplitter
    from llama_index.embeddings.ollama import OllamaEmbedding
    from llama_index.vector_stores.qdrant import QdrantVectorStore
    from qdrant_client import QdrantClient
    from qdrant_client.models import Distance, VectorParams

    st("ing", status="starting", step="connecting to Qdrant...")
    Settings.embed_model = OllamaEmbedding(model_name=ING_EMBED_MODEL)

    client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT, timeout=60)
    if not client.collection_exists(ING_COLLECTION):
        client.create_collection(ING_COLLECTION,
            vectors_config=VectorParams(size=768, distance=Distance.COSINE))

    vector_store = QdrantVectorStore(client=client, collection_name=ING_COLLECTION,
                                     batch_size=ING_BATCH_SIZE)
    pipeline = IngestionPipeline(
        transformations=[
            SentenceSplitter(chunk_size=ING_CHUNK_SIZE, chunk_overlap=ING_CHUNK_OVERLAP),
            Settings.embed_model,
        ],
        vector_store=vector_store
    )

    add_log("ING", f"Connected — {client.count(ING_COLLECTION).count:,} vectors")

    while True:
        tracker = _load_json(INGEST_TRACKER)
        pending = [fp for fp in Path(RAG_ROOT).rglob("*.summary.json")
                   if str(fp) not in tracker]

        st("ing", available=len(pending) + snap("ing")["done"],
           vectors=client.count(ING_COLLECTION).count)

        if not pending:
            st("ing", status="waiting", current="",
               step=f"waiting for new summaries ({ING_POLL_SECS}s)")
            time.sleep(ING_POLL_SECS)
            continue

        st("ing", status="running")
        add_log("ING", f"{len(pending):,} summaries to ingest")

        for i, fp in enumerate(pending, 1):
            name = fp.name
            st("ing", current=name, step="reading JSON")
            try:
                obj = json.loads(fp.read_text(encoding="utf-8"))
                st("ing", step="embedding")
                doc = Document(
                    text=_build_doc_text(obj),
                    metadata={
                        "original_path"  : obj.get("original_path", ""),
                        "rag_path"       : str(fp),
                        "file_name"      : obj.get("file_name", fp.stem),
                        "file_extension" : obj.get("file_extension", ""),
                        "document_type"  : obj.get("document_type", ""),
                        "key_topics"     : ", ".join(obj.get("key_topics", [])),
                        "entities"       : ", ".join(obj.get("entities", [])),
                        "file_size_bytes": obj.get("file_size_bytes", 0),
                        "file_modified"  : obj.get("file_modified", ""),
                    },
                    id_=obj.get("original_path") or str(fp),
                )
                nodes = pipeline.run(documents=[doc], show_progress=False)
                tracker[str(fp)] = {"status": "ok", "nodes": len(nodes), "processed_at": time.time()}
                _save_json(INGEST_TRACKER, tracker)
                with _lock:
                    _state["ing"]["done"] += 1
                    _state["ing"]["vectors"] = client.count(ING_COLLECTION).count
                add_log("ING", f"DONE {name}  ({len(nodes)} nodes)")

            except Exception as e:
                tracker[str(fp)] = {"status": "error", "error": str(e), "processed_at": time.time()}
                _save_json(INGEST_TRACKER, tracker)
                with _lock:
                    _state["ing"]["errors"] += 1
                add_log("ING", f"ERR  {name} — {type(e).__name__}")

    st("ing", status="done", current="", step="complete")


# ══════════════════════════════════════════════════════════════════════════════
# DISPLAY
# ══════════════════════════════════════════════════════════════════════════════

def _status_icon(status):
    return {"running": "● RUNNING", "scanning": "◌ SCANNING",
            "waiting": "◷ WAITING", "done": "✓ DONE",
            "starting": "… STARTING", "stopped": "✗ STOPPED"}.get(status, status.upper())

def render():
    s = snap("sum")
    g = snap("ing")
    with _lock:
        logs = list(_state["log"])

    W = 70
    os.system('cls')
    print('═' * W)
    print(f'  RAG PIPELINE   {datetime.now().strftime("%Y-%m-%d  %H:%M:%S")}')
    print('═' * W)

    # ── Summarizer ───────────────────────────────────────────────────────────
    print(f'  SUMMARIZER   {_status_icon(s["status"])}')
    print('  ' + '─' * (W - 2))
    print(f'  {(s["current"] or "—")[:60]}')
    print(f'  {s["step"][:66]}' if s["step"] else "")
    if s["total"]:
        remaining = s["total"] - s["done"] - s["errors"]
        pct_remaining = remaining / s["total"] * 100
        filled = round(46 * (100 - pct_remaining) / 100)
        print(f'  Remaining: {remaining:,}   {pct_remaining:.2f}%')
        print(f'  {"█"*filled}{"░"*(46-filled)}')
    else:
        print(f'  Remaining: (scanning...)')
    print()

    # ── Ingester ─────────────────────────────────────────────────────────────
    print(f'  INGESTER   {_status_icon(g["status"])}')
    print('  ' + '─' * (W - 2))
    print(f'  {(g["current"] or "—")[:60]}')
    print(f'  {g["step"][:66]}' if g["step"] else "")
    total = s["total"] or 0
    if total:
        remaining = total - g["done"]
        pct_remaining = remaining / total * 100
        filled = round(46 * g["done"] / total)
        print(f'  {g["done"]:,} of {total:,}   Remaining: {remaining:,}   {pct_remaining:.2f}%   Vectors: {g["vectors"]:,}')
        print(f'  {"█"*filled}{"░"*(46-filled)}')
    else:
        print(f'  Remaining: (waiting for total count...)')
    print()

    # ── Log ──────────────────────────────────────────────────────────────────
    if logs:
        print('  ' + '─' * (W - 2))
        for t, tag, msg in logs[-6:]:
            print(f'  {t}  {tag:<3}  {msg[:58]}')

    print()
    print(f'  Ctrl+C to stop both processes')
    print('═' * W)


def main():
    Path(RAG_ROOT).mkdir(parents=True, exist_ok=True)

    t_sum = threading.Thread(target=summarizer_thread, daemon=True, name="summarizer")
    t_ing = threading.Thread(target=ingester_thread,   daemon=True, name="ingester")
    t_sum.start()
    t_ing.start()

    try:
        while True:
            render()
            time.sleep(DISPLAY_REFRESH)
            # Mark stopped if thread died
            if not t_sum.is_alive() and snap("sum")["status"] not in ("done",):
                st("sum", status="stopped")
            if not t_ing.is_alive() and snap("ing")["status"] not in ("done",):
                st("ing", status="stopped")
    except KeyboardInterrupt:
        st("sum", status="stopped")
        st("ing", status="stopped")
        render()
        print("\nStopped.")


if __name__ == "__main__":
    main()
