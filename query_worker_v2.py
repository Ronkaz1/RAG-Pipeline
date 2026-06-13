import time
import traceback
import logging
import json
import pymssql
from pathlib import Path
from datetime import datetime
from openai import OpenAI
from llama_index.core import VectorStoreIndex, Settings, StorageContext
from llama_index.vector_stores.qdrant import QdrantVectorStore
from llama_index.embeddings.ollama import OllamaEmbedding
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams

# ========================= CONFIG =========================
DB_SERVER   = "YOUR_SQL_SERVER"
DB_PORT     = 1433
DB_NAME     = "ged"
DB_USER     = "YOUR_DB_USER"
DB_PASSWORD = "YOUR_DB_PASSWORD"

QDRANT_HOST      = "localhost"
QDRANT_PORT      = 6333
COLLECTION_NAME  = "rag_summaries"
EMBED_MODEL      = "nomic-embed-text"
EMBED_DIM        = 768

OLLAMA_LLM_MODEL = "qwen2.5:3b"
OLLAMA_BASE_URL  = "http://localhost:11434/v1"

GROK_API_KEY = "YOUR_GROK_API_KEY"
GROK_MODEL   = "grok-3-mini"

TOP_K            = 8     # summary nodes to retrieve
MAX_SOURCE_DOCS  = 3     # original docs to load into prompt
MAX_DOC_CHARS    = 6000  # chars per original doc (3 x 6000 = 18k fits in 32k context)
POLL_INTERVAL    = 3
# =========================================================

logging.basicConfig(
    filename="query_worker_v2.log",
    level=logging.ERROR,
    format="%(asctime)s %(levelname)s: %(message)s"
)

def ts():
    return datetime.now().strftime("%H:%M:%S")

def log(msg: str):
    print(msg, flush=True)

def log_error(msg: str, exc: BaseException = None):
    print(msg, flush=True)
    logging.error(msg)
    if exc:
        logging.error(traceback.format_exc())

ollama_client = OpenAI(api_key="ollama", base_url=OLLAMA_BASE_URL, timeout=180.0)
grok_client   = OpenAI(api_key=GROK_API_KEY, base_url="https://api.x.ai/v1", timeout=60.0)


def get_db_connection():
    return pymssql.connect(server=DB_SERVER, port=DB_PORT, database=DB_NAME,
                           user=DB_USER, password=DB_PASSWORD, conn_properties='')


def extract_text(file_path: Path, max_chars: int = MAX_DOC_CHARS) -> tuple[str, str]:
    """Return (text, description). Lightweight version for prompt loading."""
    ext = file_path.suffix.lower()
    try:
        if ext == '.pdf':
            from pypdf import PdfReader
            reader = PdfReader(str(file_path))
            text = ""
            for i, page in enumerate(reader.pages):
                text += (page.extract_text() or "") + "\n"
                if len(text) >= max_chars:
                    break
            return text[:max_chars], f"{min(i+1, len(reader.pages))}/{len(reader.pages)} pages"

        elif ext in ('.docx', '.doc'):
            from docx import Document
            doc = Document(str(file_path))
            text = "\n".join(p.text for p in doc.paragraphs)
            return text[:max_chars], f"{len(doc.paragraphs)} paragraphs"

        elif ext in ('.pptx', '.ppt'):
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                if len(text) >= max_chars:
                    break
            return text[:max_chars], f"{len(prs.slides)} slides"

        elif ext in ('.xlsx', '.xls'):
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text = ""
            for sheet in wb.worksheets:
                text += f"[Sheet: {sheet.title}]\n"
                for row in sheet.iter_rows(values_only=True):
                    row_str = "  |  ".join(str(c) for c in row if c is not None)
                    if row_str.strip():
                        text += row_str + "\n"
                    if len(text) >= max_chars:
                        break
                if len(text) >= max_chars:
                    break
            return text[:max_chars], "spreadsheet"

        else:
            text = file_path.read_text(errors='ignore')
            return text[:max_chars], f"{text.count(chr(10))} lines"

    except Exception as e:
        return "", f"load failed: {e}"


def retrieve_and_load(retriever, question: str) -> list[dict]:
    """
    Retrieve top summary nodes, deduplicate by original_path,
    then load the actual source documents from bigbertha.
    Returns list of source dicts with summary + doc text.
    """
    nodes = retriever.retrieve(question)

    # Deduplicate by original_path, keep highest-scoring node per doc
    seen = {}
    for node in nodes:
        orig = node.node.metadata.get("original_path", "")
        if not orig:
            continue
        if orig not in seen or node.score > seen[orig]["score"]:
            seen[orig] = {
                "original_path" : orig,
                "file_name"     : node.node.metadata.get("file_name", Path(orig).name),
                "document_type" : node.node.metadata.get("document_type", ""),
                "key_topics"    : node.node.metadata.get("key_topics", ""),
                "score"         : node.score,
                "summary_text"  : node.node.text,
            }

    sources = sorted(seen.values(), key=lambda x: x["score"], reverse=True)[:MAX_SOURCE_DOCS]

    log(f"\n   [{ts()}] â”€â”€ Retrieved {len(nodes)} nodes â†’ {len(seen)} unique docs â†’ loading top {len(sources)}")
    log(f"   {'â”€'*60}")

    # Load original documents
    for src in sources:
        fp = Path(src["original_path"])
        score_pct = src["score"] * 100
        if fp.exists():
            log(f"   [{ts()}] LOAD  {src['file_name']}  (score {score_pct:.1f}%)")
            text, desc = extract_text(fp)
            src["doc_text"]  = text
            src["doc_desc"]  = desc
            src["doc_loaded"] = True
            log(f"          â†’ {len(text):,} chars loaded  ({desc})")
        else:
            log(f"   [{ts()}] MISS  {src['file_name']}  â€” not reachable at {src['original_path']}")
            src["doc_text"]   = ""
            src["doc_desc"]   = "file not found"
            src["doc_loaded"] = False

    log(f"   {'â”€'*60}\n")
    return sources


def build_prompt(question: str, sources: list[dict]) -> str:
    loaded = [s for s in sources if s.get("doc_loaded")]
    summary_only = [s for s in sources if not s.get("doc_loaded")]

    parts = [
        "You are an expert analyst in the energy and oil industry. "
        "Answer the question using only the document content provided below. "
        "Be specific, cite document names where relevant, and give a complete answer.\n"
    ]

    if loaded:
        parts.append("=== FULL DOCUMENT CONTENT ===")
        for i, src in enumerate(loaded, 1):
            parts.append(
                f"\n[Document {i}: {src['file_name']}  |  "
                f"type={src['document_type']}  |  relevance={src['score']*100:.1f}%]\n"
                f"{src['doc_text']}"
            )

    if summary_only:
        parts.append("\n=== DOCUMENT SUMMARIES (originals unavailable) ===")
        for src in summary_only:
            parts.append(
                f"\n[{src['file_name']}  |  type={src['document_type']}]\n"
                f"{src['summary_text']}"
            )

    parts.append(f"\n=== QUESTION ===\n{question}")
    parts.append("\nAnswer:")
    return "\n".join(parts)


def call_ollama(prompt: str) -> str:
    response = ollama_client.chat.completions.create(
        model=OLLAMA_LLM_MODEL,
        messages=[{"role": "user", "content": prompt}]
    )
    return response.choices[0].message.content.strip()


def call_grok(question: str) -> str:
    response = grok_client.chat.completions.create(
        model=GROK_MODEL,
        messages=[
            {"role": "system", "content": (
                "You are an expert in the energy and oil industry. "
                "Answer directly without opening with industry qualifiers. "
                "Provide a complete, self-contained answer."
            )},
            {"role": "user", "content": question}
        ]
    )
    return response.choices[0].message.content.strip()


def is_empty_response(answer: str) -> bool:
    low = answer.lower()
    phrases = ["no information", "no relevant", "don't have", "do not have",
               "cannot find", "no data", "not found", "unable to find",
               "no documents", "empty", "no context"]
    return len(answer.strip()) < 80 or any(p in low for p in phrases)


def setup():
    log(f"[{ts()}] Connecting to Qdrant '{COLLECTION_NAME}'...")
    Settings.embed_model = OllamaEmbedding(model_name=EMBED_MODEL, embed_batch_size=50)

    client = QdrantClient(host=QDRANT_HOST, port=QDRANT_PORT)
    if not client.collection_exists(COLLECTION_NAME):
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(size=EMBED_DIM, distance=Distance.COSINE)
        )

    count = client.count(COLLECTION_NAME).count
    log(f"   {count:,} vectors in collection")

    vector_store     = QdrantVectorStore(client=client, collection_name=COLLECTION_NAME)
    storage_context  = StorageContext.from_defaults(vector_store=vector_store)
    index            = VectorStoreIndex.from_vector_store(vector_store, storage_context=storage_context)
    retriever        = index.as_retriever(similarity_top_k=TOP_K)
    return index, retriever


def ingest_grok_answer(index, question: str, answer: str):
    """Store the Grok Q&A pair as a vector so future similar questions hit RAG."""
    from llama_index.core import Document
    from llama_index.core.node_parser import SentenceSplitter
    doc = Document(
        text=f"Question: {question}\n\nAnswer: {answer}",
        metadata={
            "original_path" : "",
            "file_name"     : "grok_answer",
            "document_type" : "grok_qa",
            "key_topics"    : "",
            "entities"      : "",
            "source"        : "grok_fallback",
        },
        id_=f"grok_{hash(question)}",
    )
    nodes = SentenceSplitter(chunk_size=1024, chunk_overlap=64).get_nodes_from_documents([doc])
    index.insert_nodes(nodes)
    log(f"   [{ts()}] â†‘ Grok answer ingested into RAG ({len(nodes)} node(s)) â€” available for future queries")


# ========================= MAIN LOOP =========================
log(f"[{ts()}] Starting Query Worker v2  (RAG â†’ original docs â†’ LLM)")
log(f"         Collection : {COLLECTION_NAME}")
log(f"         LLM        : {OLLAMA_LLM_MODEL}  (fallback: {GROK_MODEL})")
log(f"         Top-K      : {TOP_K} nodes  â†’  top {MAX_SOURCE_DOCS} docs loaded\n")

index     = None
retriever = None

while True:
    try:
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute("""
            SELECT TOP 5 AI_Seq, Question, UserFullName
            FROM [dbo].[F_AI_Query]
            WHERE Done = 0 AND (Running = 0 OR Running IS NULL)
            ORDER BY AI_Seq ASC
        """)
        rows = cursor.fetchall()
        conn.close()

        if rows and retriever is None:
            index, retriever = setup()

        for row in rows:
            ai_seq   = row[0]
            question = row[1]
            user     = row[2] or "User"

            log(f"\n{'â•'*64}")
            log(f"[{ts()}] Q #{ai_seq}  {user}: {question[:80]}")
            log(f"{'â•'*64}")

            # Mark running
            conn = get_db_connection()
            conn.cursor().execute(
                "UPDATE [dbo].[F_AI_Query] SET Running = 1 WHERE AI_Seq = %s", (ai_seq,))
            conn.commit()
            conn.close()

            answer = ""

            sources = []
            if retriever:
                try:
                    sources = retrieve_and_load(retriever, question)

                    if sources:
                        prompt = build_prompt(question, sources)
                        log(f"   [{ts()}] Calling {OLLAMA_LLM_MODEL} with {len(prompt):,} char prompt...")
                        t0 = time.time()
                        answer = call_ollama(prompt)
                        log(f"   [{ts()}] Response in {time.time()-t0:.1f}s  ({len(answer):,} chars)")
                    else:
                        log(f"   [{ts()}] No relevant documents found in RAG.")

                except Exception as e:
                    log_error(f"   âŒ RAG error: {e}", e)
                    answer = ""

            if not answer or is_empty_response(answer):
                log(f"\n   [{ts()}] â†’ Falling back to Grok ({GROK_MODEL})...")
                sources = []   # no source docs for Grok answers
                try:
                    answer = call_grok(question)
                    log(f"   [{ts()}] Grok answered ({len(answer):,} chars)")
                    if index:
                        try:
                            ingest_grok_answer(index, question, answer)
                        except Exception as ie:
                            log_error(f"   âš  Could not ingest Grok answer: {ie}", ie)
                except Exception as e:
                    log_error(f"   âŒ Grok error: {e}", e)
                    answer = answer or f"Error: {e}"

            # Append source paths so Angular can parse them into buttons
            source_paths = [s["original_path"] for s in sources if s.get("original_path")]
            if source_paths:
                sources_block = "\n\n[SOURCES]\n" + "\n".join(source_paths) + "\n[/SOURCES]"
                answer = answer + sources_block
                log(f"   [{ts()}] Sources appended: {len(source_paths)} path(s)")
                for p in source_paths:
                    log(f"            {p}")

            log(f"\n   â”€â”€ ANSWER {'â”€'*50}")
            for line in answer[:600].splitlines():
                log(f"   {line}")
            if len(answer) > 600:
                log(f"   ... [{len(answer)-600:,} more chars]")
            log(f"   {'â”€'*58}\n")

            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                UPDATE [dbo].[F_AI_Query]
                SET Answer = %s, Raw_Answer = %s, Done = 1, Running = 0
                WHERE AI_Seq = %s
            """, (answer[:8000], answer, ai_seq))
            conn.commit()
            conn.close()

            log(f"   [{ts()}] âœ“ Saved AI_Seq {ai_seq}")

        time.sleep(POLL_INTERVAL)

    except KeyboardInterrupt:
        log("Shutting down.")
        break
    except BaseException as e:
        log_error(f"âŒ Main loop error: {type(e).__name__}: {e}", e)
        traceback.print_exc()
        err_str = str(e).lower()
        if "connection failed" in err_str or "adaptive server" in err_str:
            log("   â†» DB error â€” retrying in 30s...")
            time.sleep(30)
        else:
            time.sleep(5)

log("Worker exited.")

